from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- 디자인 테마 설정 (Dark Mode 느낌) ---
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(30, 30, 30) # Dark Grey Background

    def add_title_slide(prs, title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 69, 0) # Red Orange
        title.text_frame.paragraphs[0].font.bold = True
        
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
        return slide

    def add_content_slide(prs, title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide)
        
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 140, 0) # Lighter Orange
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        content = slide.placeholders[1]
        content.text = content_text
        
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(230, 230, 230)
            paragraph.font.size = Pt(18)
            paragraph.space_after = Pt(10)
            
        return slide

    # 1. 표지
    add_title_slide(prs, 
                    "Into the Dead: Our Darkest Days\n역기획서 (Reverse Game Design)", 
                    "작성자: AI 기획자\n주제: 생존 공포와 쉘터 관리의 메커니즘 분석")

    # 2. 개요
    add_content_slide(prs, "1. 프로젝트 개요 (Overview)",
                      "- 게임명: Into the Dead: Our Darkest Days\n"
                      "- 개발사: PikPok\n"
                      "- 장르: 2.5D 횡스크롤 쉘터 서바이벌 (Shelter Survival)\n"
                      "- 출시일: 2025년 4월 (Early Access)\n"
                      "- 배경: 1980년 텍사스 월튼 시티, 좀비 아포칼립스\n"
                      "- USP (Unique Selling Point):\n"
                      "  1) 횡스크롤 러닝 액션 IP의 장르적 변주 (액션 → 시뮬레이션)\n"
                      "  2) 심리적 압박감을 강조한 인물 관리 시스템")

    # 3. 기획 의도
    add_content_slide(prs, "2. 기획 의도 분석 (Design Intent)",
                      "■ 왜 '쉘터 생존' 인가?\n"
                      "- 기존 IP의 한계 극복: 단순 반복 플레이(Running)에서 장기적 리텐션을 위한 관리(Management) 요소 도입.\n"
                      "- 시장 포지셔닝: 'This War of Mine'의 감성 + 'State of Decay'의 액션성을 2.5D로 결합.\n\n"
                      "■ 핵심 경험 목표 (Core Experience)\n"
                      "- '무력감 속의 희망': 압도적인 좀비 수(배경)와 취약한 생존자(전경)의 대비를 통해 긴장감 극대화.\n"
                      "- 선택의 무게: 모든 파밍과 이동이 생존자의 영구적 죽음(Perma-death)과 직결되도록 설계.")

    # 4. 핵심 루프
    add_content_slide(prs, "3. 핵심 게임 루프 (Core Loop)",
                      "1. 계획 (Preparation)\n"
                      "   - 생존자 상태 확인, 물자 배분, 탐색 지역 선정\n"
                      "2. 탐색 (Scavenge - Night/Day)\n"
                      "   - 잠입(Stealth) 위주의 2.5D 횡스크롤 액션\n"
                      "   - 소음 관리(Noise Management) 및 자원 확보\n"
                      "3. 관리 (Shelter Management)\n"
                      "   - 쉘터 보강, 식량 배급, 부상 치료\n"
                      "4. 방어 및 위기 (Crisis)\n"
                      "   - 좀비 웨이브 방어, 내부 갈등 발생, 무작위 이벤트\n"
                      "➔ 최종 목표: 도시 탈출 계획 수립 및 실행")

    # 5. 상세 시스템 분석 - 2.5D 메커니즘
    add_content_slide(prs, "4. 시스템 상세: 2.5D 심도와 공포",
                      "■ 전경(Foreground) vs 후경(Background)의 분리\n"
                      "- 메커니즘: 플레이어는 전경으로만 이동 가능하지만, 위협(좀비)은 후경에서 전경으로 넘어옴.\n"
                      "- 의도: 시각적 정보는 주되 물리적 대응은 불가능한 상태를 만들어 '심리적 공포' 유발.\n\n"
                      "■ 소음과 시야 (Noise & Visibility)\n"
                      "- 시각적 요소(UI)로 소음 파동을 시각화하여 직관적인 스텔스 플레이 유도.\n"
                      "- 조명: 손전등 사용 시 시야 확보(이점) vs 좀비 어그로(리스크)의 딜레마 제공.")

    # 6. 상세 시스템 분석 - 생존자 관리
    add_content_slide(prs, "5. 시스템 상세: 욕구와 트라우마",
                      "■ Maslow의 욕구 단계 적용\n"
                      "- 생리적 욕구: 배고픔, 피로도 (기본 생존)\n"
                      "- 안전 욕구: 쉘터 내구도, 무기 상태\n"
                      "- 애정/존경 욕구: 생존자 간의 관계(유대감), 사기(Morale)\n\n"
                      "■ 트라우마 시스템 (Sanity)\n"
                      "- 단순 HP가 아닌 '정신력' 관리 필요.\n"
                      "- 동료의 죽음이나 살인은 영구적인 트라우마 특성(Trait)을 부여하여 캐릭터 성능 변화 유도.")

    # 7. UI/UX 분석
    add_content_slide(prs, "6. UI/UX 분석",
                      "■ 스큐어모피즘(Skeuomorphism)과 미니멀리즘의 조화\n"
                      "- 80년대 레트로 감성을 살린 지도, 무전기 등의 오브젝트 디자인.\n"
                      "- 인게임 HUD는 최소화하여 몰입감 증대 (체력바 대신 캐릭터 애니메이션으로 상태 표현).\n\n"
                      "■ 정보 위계 (Hierarchy)\n"
                      "- 1순위: 소음 지표 (생존 직결)\n"
                      "- 2순위: 상호작용 아이콘 (파밍)\n"
                      "- 3순위: 캐릭터 상태창 (메뉴 진입 시 확인)")

    # 8. SWOT 분석
    add_content_slide(prs, "7. SWOT 분석",
                      "■ Strengths (강점)\n"
                      "- 검증된 'Into the Dead' IP 인지도.\n"
                      "- 2.5D 그래픽의 고유한 아트 스타일과 연출.\n\n"
                      "■ Weaknesses (약점)\n"
                      "- 반복적인 파밍 구조로 인한 중후반 지루함.\n"
                      "- 횡스크롤의 한계로 인한 전투 전략의 단순함.\n\n"
                      "■ Opportunities (기회)\n"
                      "- 스트리밍 친화적 게임성 (긴장감, 점프 스케어).\n"
                      "- 샌드박스 모드 추가를 통한 리플레이성 강화.\n\n"
                      "■ Threats (위협)\n"
                      "- 'Project Zomboid' 등 강력한 경쟁작의 존재.\n"
                      "- 얼리 액세스 기간 중 콘텐츠 고갈 우려.")

    # 9. 개선 제안
    add_content_slide(prs, "8. 역기획 제안: '환각(Hallucination)' 시스템",
                      "■ 문제 정의\n"
                      "- 후반부 자원이 풍족해지면 좀비에 대한 공포감이 '귀찮음'으로 변질됨.\n\n"
                      "■ 제안 내용: [공포의 실체화]\n"
                      "- 정신력(Sanity)이 낮아지면 2.5D 배경의 좀비가 전경으로 튀어나오는 '환각' 발생.\n"
                      "- 실제 데미지는 없으나, 총기 발사 유도(탄약 낭비) 및 소음 유발로 실제 좀비를 불러모음.\n\n"
                      "■ 기대 효과\n"
                      "- 자원이 많아도 정신력 관리가 안 되면 게임이 터질 수 있는 변수 창출.\n"
                      "- 후반부 긴장감 유지 및 내러티브 몰입도 강화.")

    # 10. 결론
    add_content_slide(prs, "9. 결론 (Conclusion)",
                      "- 'Into the Dead: Our Darkest Days'는 액션에서 시뮬레이션으로의 성공적인 장르 확장을 보여줌.\n"
                      "- 2.5D의 공간적 제약을 시스템적 공포로 승화시킨 점이 백미.\n"
                      "- 제안된 '환각 시스템' 도입 시, 심리적 공포 게임으로서의 정체성을 더욱 확고히 할 수 있을 것임.")

    # 저장
    prs.save('Into_The_Dead_Reverse_Plan.pptx')
    print("PPT 파일이 생성되었습니다: Into_The_Dead_Reverse_Plan.pptx")

if __name__ == "__main__":
    create_presentation()